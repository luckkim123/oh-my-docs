"""Canonical pptx shape-property assertion — the literal file doc-builder.md step 6 names
(`.omd/<slug>/assert_shapes.py`). Data file, copy/adapt into your own per-job build script
(see references/snippets/README.md) — never imported by an agent at runtime, only by this
repo's own test suite.

Zero-dep import: `pptx` is only required when `assert_shapes()` is actually called, never at
module load — `from __future__ import annotations` (PEP 563) defers every annotation to an
unevaluated string, and `pptx` is imported only under `TYPE_CHECKING`, so this module collects
and its no-pptx tests run on a machine without python-pptx installed.

Referenced by: agents/doc-builder.md step 6 (mandatory self-gate) and agents/doc-verifier.md
step 2b (independent re-run) — both point at this same function so the two checks cannot
silently diverge into two hand-written implementations of "the same four checks".
"""
from __future__ import annotations

from typing import TYPE_CHECKING

if TYPE_CHECKING:
    import pptx

_PLACEHOLDER_MARKERS = ("click to add", "[insert", "todo", "lorem")


def has_placeholder_text(s: str) -> bool:
    """Pure substring check, factored out so it is testable without python-pptx installed."""
    low = (s or "").lower()
    return any(marker in low for marker in _PLACEHOLDER_MARKERS)


def assert_shapes(
    prs: "pptx.Presentation",
    expected_font_name: str,
    outline_slide_count: int | None = None,
) -> list:
    """Violation strings; empty list = ASSERT OK.
    Checks every body run's `font.size is not None`; every shape's `width>0 and height>0`;
    `run.font.name == expected_font_name`; `has_placeholder_text` on every text frame;
    optional slide-count check against the outline."""
    violations = []
    slides = list(prs.slides)
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.width <= 0 or shape.height <= 0:
                violations.append(f"slide {i}: shape {shape.shape_id!r} has width/height <= 0")
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if has_placeholder_text(tf.text):
                violations.append(f"slide {i}: shape {shape.shape_id!r} still shows placeholder text")
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        violations.append(
                            f"slide {i}: shape {shape.shape_id!r} run {run.text!r} has no explicit font.size"
                        )
                    if run.font.name != expected_font_name:
                        violations.append(
                            f"slide {i}: shape {shape.shape_id!r} run {run.text!r} "
                            f"font.name={run.font.name!r} != {expected_font_name!r}"
                        )
    if outline_slide_count is not None and len(slides) != outline_slide_count:
        violations.append(f"slide count {len(slides)} != outline's {outline_slide_count}")
    return violations


def main() -> int:
    import sys

    if len(sys.argv) < 3:
        print("usage: assert_shapes.py <deck.pptx> <expected-font> [expected-count]", file=sys.stderr)
        return 2
    import pptx as _pptx

    deck_path, expected_font = sys.argv[1], sys.argv[2]
    expected_count = int(sys.argv[3]) if len(sys.argv) > 3 else None
    prs = _pptx.Presentation(deck_path)
    violations = assert_shapes(prs, expected_font, expected_count)
    if violations:
        for v in violations:
            print(v)
        return 1
    print("ASSERT OK")
    return 0


if __name__ == "__main__":
    import sys

    sys.exit(main())
