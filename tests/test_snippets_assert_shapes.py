"""references/snippets/assert_shapes.py — pptx shape-property assertion (the doc-builder
step 6 / doc-verifier step 2b canonical check). has_placeholder_text is pure Python (no
skip); the module must import cleanly with pptx absent (TYPE_CHECKING guard proof);
assert_shapes() itself needs a real pptx.Presentation, skipped when python-pptx absent.
"""
import importlib.util
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT / "references" / "snippets"))
import assert_shapes as az  # noqa: E402 — this import itself proves collection works w/o pptx


def test_module_imports_cleanly_without_pptx_installed():
    # If this file collected at all (see the top-level import above), the TYPE_CHECKING
    # guard already did its job on a machine without python-pptx. This test just names
    # that fact explicitly for the reader.
    assert az is not None


@pytest.mark.parametrize(
    "text,expected",
    [
        ("Click to add text", True),
        ("[Insert chart]", True),
        ("TODO", True),
        ("Lorem ipsum dolor sit amet", True),
        ("Q3 revenue grew 12% year over year.", False),
        ("", False),
    ],
)
def test_has_placeholder_text(text, expected):
    assert az.has_placeholder_text(text) is expected


@pytest.mark.skipif(
    importlib.util.find_spec("pptx") is None, reason="python-pptx not installed"
)
def test_assert_shapes_callable_with_real_pptx():
    import pptx

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    box.text_frame.text = "hello"
    box.text_frame.paragraphs[0].runs[0].font.size = pptx.util.Pt(18)
    box.text_frame.paragraphs[0].runs[0].font.name = "Arial"
    violations = az.assert_shapes(prs, expected_font_name="Arial")
    assert violations == []
